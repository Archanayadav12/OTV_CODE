import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# Load the dataset
file_path = r"C:\Users\yarch\OneDrive\Desktop\OTV_CODE\Dataset -2_new.xlsx"  # Use raw string to avoid unicode error
data = pd.ExcelFile(file_path)
df = data.parse('Dataset -2')

# Refund Rates Analysis
refund_rate = df.groupby('payment_method')['refund_status'].apply(lambda x: (x == 'Yes').mean()).sort_values()

# Visualize Refund Rates
plt.figure(figsize=(10, 6))
refund_rate.plot(kind='bar')
plt.title('Refund Rates by Payment Method')
plt.ylabel('Refund Rate')
plt.xlabel('Payment Method')
img_path_refund = r"C:\Users\yarch\OneDrive\Desktop\OTV_CODE\refund_rates_updated.png"  # Save image locally
plt.savefig(img_path_refund)
plt.close()

# Top Suppliers by Revenue
supplier_revenue = df.groupby('supplier_id')['selling_price'].sum().sort_values(ascending=False)

# Visualize Top Suppliers
plt.figure(figsize=(10, 6))
supplier_revenue.head(10).plot(kind='bar')
plt.title('Top 10 Suppliers by Revenue')
plt.ylabel('Total Revenue')
plt.xlabel('Supplier ID')
img_path_supplier = r"C:\Users\yarch\OneDrive\Desktop\OTV_CODE\supplier_revenue_updated.png"  # Save image locally
plt.savefig(img_path_supplier)
plt.close()

# Monthly Sales Trends
sales_trend = df.groupby(df['booking_date'].dt.to_period('M'))['selling_price'].sum()

# Visualize Sales Trends
plt.figure(figsize=(12, 6))
sales_trend.plot()
plt.title('Monthly Sales Trend')
plt.ylabel('Total Sales')
plt.xlabel('Booking Month')
img_path_sales = r"C:\Users\yarch\OneDrive\Desktop\OTV_CODE\sales_trends_updated.png"  # Save image locally
plt.savefig(img_path_sales)
plt.close()

# Create a PowerPoint Presentation
presentation = Presentation()

# Add the refund rate image to the PowerPoint
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide.shapes.title
title.text = "Refund Rates by Payment Method"
slide.shapes.add_picture(img_path_refund, Inches(1), Inches(1), width=Inches(8.5))

# Add the supplier revenue image to the PowerPoint
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide.shapes.title
title.text = "Top 10 Suppliers by Revenue"
slide.shapes.add_picture(img_path_supplier, Inches(1), Inches(1), width=Inches(8.5))

# Add the sales trend image to the PowerPoint
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
title = slide.shapes.title
title.text = "Monthly Sales Trend"
slide.shapes.add_picture(img_path_sales, Inches(1), Inches(1), width=Inches(8.5))

# Save the PowerPoint presentation
pptx_path = r"C:\Users\yarch\OneDrive\Desktop\OTV_CODE\OTV_Analysis_Presentation.pptx"  # Save PowerPoint locally
presentation.save(pptx_path)

pptx_path  # Return the path to the saved PowerPoint file
